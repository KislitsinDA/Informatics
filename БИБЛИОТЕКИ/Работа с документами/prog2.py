from pptx import Presentation

prs = Presentation()

slide_layout = prs.slide_layouts[1]

slide_1 = prs.slides.add_slide(slide_layout)
slide_1.shapes.title.text = "Random"
slide_1.placeholders[1].text = "random()  method of ransom.Random instance\nrandom() -> x in the interval (0, 1)"

slide_2 = prs.slides.add_slide(slide_layout)
slide_2.shapes.title.text = "Choice"
slide_2.placeholders[1].text = "choise(seq) method of random.Random instance," \
                                "\nChoose a random element from a non-empty sequence."

slide_3 = prs.slides.add_slide(slide_layout)
slide_3.shapes.title.text = "Choices"
slide_3.placeholders[1].text = "choises(population, weights=None, *, cum_weights=None, k=1" \
                                "method of random.Random instance.\nReturn a k sized list of population elements " \
                                "chosen with replacement.\nIf the relative weights cumulative weights are not " \
                                "specified, the selections are made with equal probably."

slide_4 = prs.slides.add_slide(slide_layout)
slide_4.shapes.title.text = "Randrange"
slide_4.placeholders[1].text = "randrange(start, stop=None, step=1, _int=<class 'int'>) method of random.Random" \
                                "instance.\nChoose a random item from range(start, stop[, step]).\nThis fixes the " \ 
                                "problem with randint() which includes the endpoint; in Python this is usually not " \
                                "what you want."

slide_5 = prs.slides.add_slide(slide_layout)
slide_5.shapes.title.text = "Betavariate"
slide_5.placeholders[1].text = "betavariate (alpha, beta) method of random.Random instance.\nConditions on the " \
                                "parameters are alpha > 0 and beta > 0."

prs.save("random_info.pptx")
